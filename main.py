# To avoid attribute errors if you are running a Python version above 3.9.6
import collections.abc
c = collections
c.abc = collections.abc

import tkinter.messagebox
import os
import tkinter as tk
from tkinter import ttk
from tkinter import filedialog
from pptx import Presentation
from pptx.util import Inches



def choose_input_directory():
    input_directory = filedialog.askdirectory(title='Select Input Directory')
    input_directory_label.config(text=input_directory)


def choose_output_directory():
    output_directory = filedialog.askdirectory(title='Select Output Directory')
    output_directory_label.config(text=output_directory)


def convert_images_to_pptx():
    input_directory = input_directory_label['text']
    output_directory = output_directory_label['text']
    create_new_folder = create_new_folder_var.get()

    if create_new_folder:
        output_directory = os.path.join(output_directory, 'converted_pptx')
        os.makedirs(output_directory, exist_ok=True)

    dir_list = os.listdir(input_directory)

# Get the number of files
    a = len(dir_list)

# Convert png to slides
    for filename in dir_list:
        if filename.endswith('.png') or filename.endswith('.jpg') or filename.endswith('.jpeg'):
            img_path = os.path.join(input_directory, filename)
            prs = Presentation()
            slide = prs.slides.add_slide(prs.slide_layouts[6])

            left = top = Inches(1)
            height = Inches(3.5)
            pic = slide.shapes.add_picture(img_path, left, top, height=height)

            # Check if the file already exists
            counter = 0
            new_filename = f"{os.path.splitext(filename)[0]}.pptx"
            while os.path.exists(os.path.join(output_directory, new_filename)):
                counter += 1
                new_filename = f"{os.path.splitext(filename)[0]}_{counter}.pptx"

            # Save the presentation in the specified output directory
            prs.save(os.path.join(output_directory, f'{filename}.pptx'))

    # Show a message box to notify the user that the conversion is finished
    tkinter.messagebox.showinfo(
        "Conversion finished", "The conversion is finished")


root = tk.Tk()
root.title('Convert Images to PowerPoint')
root.geometry("400x400")
root.config(bg="#3F3F3F")

# Create a custom style
style = ttk.Style()
style.configure("BW.TLabel", foreground="white", background="#3F3F3F", font=("Helvetica", 16),
                padding=20, relief="flat", borderwidth=0)
style.configure("BG.TButton", font=("Courier", 16), bg="#1F7A8C", fg="white", relief="flat",
                activebackground="#19647E", activeforeground="white")

# Add a frame for the widgets
frame = tk.Frame(root, bg="#3F3F3F", bd=0, relief="flat")
frame.pack(fill="both", expand=True, padx=20, pady=20)

input_directory_label = ttk.Label(
    frame, text='Input Directory: Not Selected', style="BW.TLabel", anchor="w")
input_directory_label.pack(pady=10)

input_directory_button = tk.Button(
    frame, text='Select Input Directory',  command=choose_input_directory,
    relief="flat", borderwidth=0, highlightthickness=0, bd=0)
input_directory_button.pack(pady=10)

output_directory_label = ttk.Label(
    frame, text='Output Directory: Not Selected', style="BW.TLabel", anchor="w")
output_directory_label.pack(pady=10)

output_directory_button = tk.Button(
    frame, text='Select Output Directory',  command=choose_output_directory,
    relief="flat", borderwidth=0, highlightthickness=0, bd=0)
output_directory_button.pack(pady=10)

create_new_folder_var = tk.IntVar()
create_new_folder_checkbox = tk.Checkbutton(
    frame, text='Create New Folder', variable=create_new_folder_var, bg="#3F3F3F", fg="white",
    activebackground="#3F3F3F", activeforeground="white", highlightthickness=0, bd=0)
create_new_folder_checkbox.pack(padx=10, pady=10)

convert_button = tk.Button(
    frame, text='Convert',  command=convert_images_to_pptx,
    relief="flat", borderwidth=0, highlightthickness=0, bd=0)
convert_button.pack(padx=10, pady=10)

root.mainloop()
